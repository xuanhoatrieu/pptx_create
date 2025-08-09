import os
print("Loading pptx_creator.py") # Debugging line
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import nsdecls
from lxml import etree

def add_background(slide, background_image_path):
    """Adds a background image to the slide."""
    left = top = Inches(0)
    pic = slide.shapes.add_picture(background_image_path, left, top, width=Inches(10), height=Inches(5.625))
    slide.shapes._spTree.remove(pic._element)
    slide.shapes._spTree.insert(2, pic._element)

def get_mime_type(filepath):
    """Determines the MIME type for common audio formats."""
    if filepath.lower().endswith('.mp3'):
        return 'audio/mpeg'
    elif filepath.lower().endswith('.wav'):
        return 'audio/wav'
    # Add other audio formats as needed
    return 'audio/mpeg' # Default fallback

def create_presentation(output_path, slides_data, title_bg, content_bg):
    """Creates a PowerPoint presentation from the processed slide data."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    # --- Title Slide ---
    if slides_data:
        title_slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(title_slide_layout)
        add_background(slide, title_bg)

        # Title
        title_shape = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
        title_frame = title_shape.text_frame
        p = title_frame.paragraphs[0]
        p.text = slides_data[0].get("title", "")
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

        # Subtitle
        if slides_data[0].get("subtitle"):
            subtitle_shape = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(1))
            subtitle_frame = subtitle_shape.text_frame
            p_sub = subtitle_frame.paragraphs[0]
            p_sub.text = slides_data[0].get("subtitle")
            p_sub.font.size = Pt(22)
            p_sub.font.color.rgb = RGBColor(255, 255, 255)
            p_sub.alignment = PP_ALIGN.CENTER

    # --- Content Slides ---
    for i, slide_info in enumerate(slides_data):
        # Special handling for the first slide's audio
        if i == 0:
            if slide_info.get("audio_path"):
                try:
                    audio_path = slide_info["audio_path"]
                    mime_type = get_mime_type(audio_path)
                    movie = prs.slides[0].shapes.add_movie(
                        audio_path, Inches(10), Inches(0), Inches(0.5), Inches(0.5), mime_type=mime_type
                    )
                    tree = movie.element.xpath('.//p:videoFile')[0]
                    tree.addprevious(etree.fromstring(f'<p:cTn {nsdecls("p")} id="1"><p:stCondLst><p:cond delay="0"/></p:stCondLst></p:cTn>'))
                except Exception as e:
                    print(f"Could not add audio to title slide {audio_path}: {e}")
            # If there's only one slide, we're done after this.
            if len(slides_data) == 1:
                break
            else:
                continue

        content_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(content_layout)
        add_background(slide, content_bg)

        # Title, Content, Image, Notes...
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
        p = title_shape.text_frame.paragraphs[0]
        p.text = slide_info.get("title", "")
        p.font.size = Pt(28)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

        if slide_info.get("image_path"):
            content_shape = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.5), Inches(4.0))
            try:
                slide.shapes.add_picture(slide_info["image_path"], Inches(5.5), Inches(1.5), height=Inches(3))
            except Exception as e:
                print(f"Could not add picture {slide_info['image_path']}: {e}")
        else:
            content_shape = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(4.0))
        
        content_frame = content_shape.text_frame
        content_frame.word_wrap = True

        for bullet in slide_info.get("bullets", []):
            p_point = content_frame.add_paragraph()
            if bullet.get("point"):
                p_point.text = f'{bullet.get("emoji", "")} {bullet.get("point")}'
                p_point.font.bold = True
                p_point.font.size = Pt(22)
                p_point.space_after = Pt(2)
                p_desc = content_frame.add_paragraph()
                p_desc.text = bullet["description"]
                p_desc.font.size = Pt(18)
                p_desc.space_before = Pt(0)
                p_desc.space_after = Pt(8)
                p_desc.level = 1
            else:
                p_point.text = bullet["description"]
                p_point.font.size = Pt(20)
            p_point.font.color.rgb = RGBColor(58, 102, 77)

        if slide_info.get("notes"):
            slide.notes_slide.notes_text_frame.text = slide_info["notes"]

        if slide_info.get("audio_path"):
            try:
                audio_path = slide_info["audio_path"]
                mime_type = get_mime_type(audio_path)
                movie = slide.shapes.add_movie(
                    audio_path, Inches(10), Inches(0), Inches(0.5), Inches(0.5), mime_type=mime_type
                )
                tree = movie.element.xpath('.//p:videoFile')[0]
                tree.addprevious(etree.fromstring(f'<p:cTn {nsdecls("p")} id="{i+2}"><p:stCondLst><p:cond delay="0"/></p:stCondLst></p:cTn>'))
            except Exception as e:
                print(f"Could not add audio {audio_path}: {e}")

    prs.save(output_path)
